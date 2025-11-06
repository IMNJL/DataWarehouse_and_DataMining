"""Generate a PPTX from `slides_text.md` and images in the repository outputs folder.

Usage:
    python src/make_presentation.py

The script looks for `slides_text.md` in the project root, parses simple slide blocks,
adds title + bullets + any referenced images, and writes `presentation.pptx` to the project root.
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDES_MD = ROOT / "slides_text.md"
OUT_PPTX = ROOT / "presentation.pptx"


def find_outputs_dir(start_dir: Path) -> Path:
    # Look for a directory named 'outputs' walking up from start_dir and then root
    cur = start_dir
    for _ in range(5):
        candidate = cur / "outputs"
        if candidate.exists() and candidate.is_dir():
            return candidate
        if cur.parent == cur:
            break
        cur = cur.parent
    # fallback to repository root outputs
    fallback = ROOT.parent / "outputs"
    return fallback if fallback.exists() else (ROOT / "outputs")


def parse_slides(md_text: str):
    # Split on lines starting with 'Slide ' (e.g., 'Slide 1 — Title')
    parts = re.split(r"(?m)^Slide\s+\d+\b[\s\S]*?\n", md_text)
    # The above leaves an initial header before first Slide; instead we'll build blocks manually
    slides = []
    lines = md_text.splitlines()
    cur = None
    for line in lines:
        m = re.match(r"^Slide\s+(\d+)\s+—\s*(.*)$", line)
        if m:
            if cur:
                slides.append(cur)
            cur = {"title": m.group(2).strip(), "bullets": [], "images": [], "notes": []}
            continue
        if cur is None:
            continue
        # detect bullets
        bm = re.match(r"^\-\s+(.*)$", line)
        if bm:
            cur["bullets"].append(bm.group(1).strip())
            continue
        # detect include figure lines: Include figure: `outputs/filename.png`
        im = re.search(r"Include figure:\s*`([^`]+)`", line)
        if im:
            cur["images"].append(im.group(1).strip())
            continue
        # speaker notes marker
        if line.strip().startswith("Speaker notes:"):
            # subsequent lines until next blank or next Slide belong to notes
            cur["notes"].append(line.replace("Speaker notes:", "").strip())
            continue
        if line.strip():
            # append remaining text to notes
            cur["notes"].append(line.strip())

    if cur:
        slides.append(cur)
    return slides


def make_presentation():
    if not SLIDES_MD.exists():
        print(f"slides_text.md not found at {SLIDES_MD}")
        return
    md = SLIDES_MD.read_text(encoding="utf8")
    slides = parse_slides(md)

    prs = Presentation()
    # Title slide from first slide block if available
    if slides:
        title_slide_layout = prs.slide_layouts[0]
        s0 = prs.slides.add_slide(title_slide_layout)
        s0.shapes.title.text = slides[0]["title"] if slides[0]["title"] else "Presentation"
        if slides[0]["bullets"]:
            try:
                s0.placeholders[1].text = "\n".join(slides[0]["bullets"])[:800]
            except Exception:
                pass

    outputs_dir = find_outputs_dir(ROOT)

    for block in slides[1:]:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = block.get("title", "")
        # add bullets
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(block.get("bullets", [])):
            p = body.add_paragraph() if i > 0 else body.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

        # add images if present
        for img_ref in block.get("images", []):
            # resolve image path: try relative to outputs_dir or directly as path
            img_path = Path(img_ref)
            if not img_path.exists():
                candidate = outputs_dir / img_ref.split('/')[-1]
                if candidate.exists():
                    img_path = candidate
                else:
                    # try root/outputs/img_ref
                    candidate2 = ROOT.parent / 'outputs' / img_ref.split('/')[-1]
                    if candidate2.exists():
                        img_path = candidate2
            if img_path.exists():
                # place image below content
                left = Inches(1)
                top = Inches(2.0)
                width = Inches(6)
                slide.shapes.add_picture(str(img_path), left, top, width=width)
            else:
                # If image not found, add a note
                note_tf = slide.notes_slide.notes_text_frame
                note_tf.text = note_tf.text + f"\nImage not found: {img_ref}"

    prs.save(OUT_PPTX)
    print(f"Saved presentation to {OUT_PPTX}")


if __name__ == '__main__':
    make_presentation()
