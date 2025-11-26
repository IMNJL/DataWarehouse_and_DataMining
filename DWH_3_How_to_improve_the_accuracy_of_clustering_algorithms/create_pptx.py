"""Generate a PPTX presentation summarizing the DPC experiments.

Requires: python-pptx, pillow

Usage:
    python create_pptx.py --images-dir outputs --slide-json slide_contents.json --out DPC_experiment_presentation.pptx
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_image_slide(prs, title, img_path, notes=None):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.3)
    slide.shapes.add_picture(img_path, left, top, width=Inches(9))
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_pptx(slide_json_path, images_dir, outpath):
    with open(slide_json_path, 'r', encoding='utf-8') as fh:
        slides = json.load(fh)

    prs = Presentation()
    for s in slides:
        typ = s.get('type', 'bullet')
        if typ == 'title':
            add_title_slide(prs, s.get('title', ''), s.get('subtitle', ''))
        elif typ == 'bullet':
            add_bullet_slide(prs, s.get('title', ''), s.get('bullets', []), notes=s.get('notes'))
        elif typ == 'image':
            img = s.get('image')
            img_path = os.path.join(images_dir, img) if not os.path.isabs(img) else img
            if os.path.exists(img_path):
                add_image_slide(prs, s.get('title', ''), img_path, notes=s.get('notes'))
            else:
                add_bullet_slide(prs, s.get('title', ''), [f'Image missing: {img_path}'])

    prs.save(outpath)
    print('Saved PPTX to', outpath)


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument('--images-dir', default='outputs')
    parser.add_argument('--slide-json', default='slide_contents.json')
    parser.add_argument('--out', default='DPC_experiment_presentation.pptx')
    args = parser.parse_args()
    build_pptx(args.slide_json, args.images_dir, args.out)
